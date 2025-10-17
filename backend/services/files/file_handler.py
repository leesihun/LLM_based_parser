#!/usr/bin/env python3
"""
File Handler Module
Handles file uploads, reading, and processing for the LLM system
"""

import os
import mimetypes
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List
import hashlib
from datetime import datetime
import json

class FileHandler:
    """Handles file operations for the LLM system"""
    
    def __init__(self, upload_dir="uploads", max_file_size=10*1024*1024):  # 10MB default
        """
        Initialize file handler
        
        Args:
            upload_dir (str): Directory for uploaded files
            max_file_size (int): Maximum file size in bytes
        """
        self.upload_dir = Path(upload_dir)
        self.max_file_size = max_file_size
        self.logger = logging.getLogger(__name__)
        
        # Create upload directory if it doesn't exist
        self.upload_dir.mkdir(exist_ok=True)
        
        # Supported file types (expanded for enhanced processing)
        self.supported_types = {
            # Text files
            '.txt': 'text',
            '.md': 'text',
            '.rst': 'text',
            '.log': 'text',
            
            # Code files
            '.py': 'code',
            '.js': 'code',
            '.ts': 'code',
            '.jsx': 'code',
            '.tsx': 'code',
            '.java': 'code',
            '.cpp': 'code',
            '.c': 'code',
            '.h': 'code',
            '.cs': 'code',
            '.php': 'code',
            '.rb': 'code',
            '.go': 'code',
            '.rs': 'code',
            '.swift': 'code',
            '.kt': 'code',
            '.scala': 'code',
            '.sql': 'code',
            '.html': 'code',
            '.css': 'code',
            '.json': 'code',
            '.xml': 'code',
            '.yaml': 'code',
            '.yml': 'code',
            '.toml': 'code',
            '.ini': 'code',
            '.conf': 'code',
            
            # Documents
            '.pdf': 'document',
            '.docx': 'document',
            '.doc': 'document',
            '.pptx': 'presentation',
            '.ppt': 'presentation',
            '.xlsx': 'spreadsheet',
            '.xls': 'spreadsheet',
            '.csv': 'spreadsheet',
            
            # Images
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.bmp': 'image',
            '.tiff': 'image',
            '.tif': 'image',
            '.webp': 'image',
            
            # Archives
            '.zip': 'archive',
            '.tar': 'archive',
            '.gz': 'archive',
            '.rar': 'archive',
            '.7z': 'archive'
        }
    
    def validate_file(self, filename: str, file_size: int) -> Dict[str, Any]:
        """
        Validate uploaded file
        
        Args:
            filename (str): Original filename
            file_size (int): File size in bytes
            
        Returns:
            Dict: Validation result
        """
        result = {
            "valid": False,
            "error": None,
            "file_type": None,
            "category": None
        }
        
        # Check file size
        if file_size > self.max_file_size:
            result["error"] = f"File too large. Maximum size: {self.max_file_size / (1024*1024):.1f}MB"
            return result
        
        # Check file extension
        file_ext = Path(filename).suffix.lower()
        if file_ext not in self.supported_types:
            result["error"] = f"Unsupported file type: {file_ext}"
            return result
        
        result["valid"] = True
        result["file_type"] = file_ext
        result["category"] = self.supported_types[file_ext]
        
        return result
    
    def generate_safe_filename(self, filename: str) -> str:
        """Generate safe filename with timestamp and hash"""
        original_name = Path(filename).stem
        extension = Path(filename).suffix
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Create hash from original filename for uniqueness
        name_hash = hashlib.md5(filename.encode()).hexdigest()[:8]
        
        # Clean original name (remove special characters)
        safe_name = "".join(c for c in original_name if c.isalnum() or c in "._-")[:50]
        
        return f"{safe_name}_{timestamp}_{name_hash}{extension}"
    
    def save_uploaded_file(self, file_data: bytes, filename: str, user_id: str) -> Dict[str, Any]:
        """
        Save uploaded file and return file info
        
        Args:
            file_data (bytes): File data
            filename (str): Original filename
            user_id (str): User ID for organization
            
        Returns:
            Dict: File save result
        """
        try:
            # Validate file
            validation = self.validate_file(filename, len(file_data))
            if not validation["valid"]:
                return {"success": False, "error": validation["error"]}
            
            # Generate safe filename
            safe_filename = self.generate_safe_filename(filename)
            
            # Create user directory
            user_dir = self.upload_dir / user_id
            user_dir.mkdir(exist_ok=True)
            
            # Save file
            file_path = user_dir / safe_filename
            with open(file_path, 'wb') as f:
                f.write(file_data)
            
            # Create file metadata
            file_info = {
                "success": True,
                "file_id": safe_filename,
                "original_name": filename,
                "safe_filename": safe_filename,
                "file_path": str(file_path),
                "file_size": len(file_data),
                "file_type": validation["file_type"],
                "category": validation["category"],
                "upload_time": datetime.now().isoformat(),
                "user_id": user_id
            }
            
            # Save metadata
            self._save_file_metadata(file_info)
            
            self.logger.info(f"File saved: {safe_filename} for user {user_id}")
            return file_info
            
        except Exception as e:
            self.logger.error(f"Error saving file {filename}: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def _save_file_metadata(self, file_info: Dict[str, Any]):
        """Save file metadata to JSON"""
        metadata_file = self.upload_dir / "file_metadata.json"
        
        # Load existing metadata
        try:
            if metadata_file.exists():
                with open(metadata_file, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)
            else:
                metadata = {}
        except:
            metadata = {}
        
        # Add new file info
        metadata[file_info["file_id"]] = file_info
        
        # Save updated metadata
        with open(metadata_file, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False)
    
    def read_file_content(self, file_id: str, user_id: str) -> Dict[str, Any]:
        """
        Read file content for processing
        
        Args:
            file_id (str): File ID
            user_id (str): User ID for access control
            
        Returns:
            Dict: File content and info
        """
        try:
            # Get file metadata
            file_info = self.get_file_info(file_id, user_id)
            if not file_info:
                return {"success": False, "error": "File not found"}
            
            file_path = Path(file_info["file_path"])
            if not file_path.exists():
                return {"success": False, "error": "File not found on disk"}
            
            content = ""
            category = file_info["category"]
            
            # Read based on file category
            if category in ["text", "code"]:
                content = self._read_text_file(file_path)
            elif category == "document":
                content = self._read_pdf_file(file_path) if file_info["file_type"] == ".pdf" else self._read_document_file(file_path)
            elif category == "presentation":
                content = self._read_presentation_file(file_path)
            elif category == "spreadsheet":
                content = self._read_spreadsheet_file(file_path)
            elif category == "image":
                content = self._read_image_file(file_path)
            elif category == "archive":
                content = self._read_archive_file(file_path)
            else:
                return {"success": False, "error": "Unsupported file category"}

            # Ensure content is not None
            if content is None:
                content = ""

            return {
                "success": True,
                "content": content,
                "file_info": file_info,
                "content_length": len(content)
            }
            
        except Exception as e:
            self.logger.error(f"Error reading file {file_id}: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def _read_text_file(self, file_path: Path) -> str:
        """Read text-based files with JSON enhancement"""
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']

        content = None
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue

        # If all encodings fail, read as binary and decode with errors ignored
        if content is None:
            with open(file_path, 'rb') as f:
                content = f.read().decode('utf-8', errors='ignore')

        # Enhanced JSON processing
        if file_path.suffix.lower() == '.json':
            return self._format_json_content(content, file_path.name)

        return content

    def _format_json_content(self, content: str, filename: str) -> str:
        """Format JSON content for better LLM comprehension"""
        try:
            # Parse JSON
            data = json.loads(content)

            # Generate formatted output with statistics
            formatted = f"JSON File: {filename}\n"
            formatted += "=" * 60 + "\n\n"

            # Add statistics
            stats = self._get_json_statistics(data)
            formatted += "STATISTICS:\n"
            formatted += f"  - Type: {stats['type']}\n"
            formatted += f"  - Size: {stats['size']}\n"
            if stats['type'] == 'array':
                formatted += f"  - Array Length: {stats['array_length']}\n"
            elif stats['type'] == 'object':
                formatted += f"  - Keys: {stats['key_count']}\n"
                formatted += f"  - Top-level fields: {', '.join(stats['keys'][:10])}\n"
                if len(stats['keys']) > 10:
                    formatted += f"    ... and {len(stats['keys']) - 10} more\n"

            formatted += f"  - Max Depth: {stats['max_depth']}\n"
            formatted += "\n" + "-" * 60 + "\n\n"

            # Add prettified JSON
            formatted += "FORMATTED CONTENT:\n\n"
            formatted += json.dumps(data, indent=2, ensure_ascii=False)

            return formatted

        except json.JSONDecodeError as e:
            # Return original content with error note if JSON is malformed
            return f"JSON File: {filename}\n" + \
                   f"WARNING: Malformed JSON - {str(e)}\n" + \
                   "=" * 60 + "\n\n" + \
                   "RAW CONTENT:\n\n" + content
        except Exception as e:
            # Fallback to original content
            self.logger.warning(f"Error formatting JSON {filename}: {str(e)}")
            return content

    def _get_json_statistics(self, data: Any) -> Dict[str, Any]:
        """Get statistics about JSON structure"""
        stats = {
            'type': type(data).__name__,
            'size': len(json.dumps(data)),
            'max_depth': self._calculate_json_depth(data),
            'array_length': 0,
            'key_count': 0,
            'keys': []
        }

        if isinstance(data, list):
            stats['type'] = 'array'
            stats['array_length'] = len(data)
        elif isinstance(data, dict):
            stats['type'] = 'object'
            stats['key_count'] = len(data)
            stats['keys'] = list(data.keys())

        return stats

    def _calculate_json_depth(self, data: Any, current_depth: int = 0) -> int:
        """Calculate maximum nesting depth of JSON structure"""
        if isinstance(data, dict):
            if not data:
                return current_depth
            return max(self._calculate_json_depth(v, current_depth + 1) for v in data.values())
        elif isinstance(data, list):
            if not data:
                return current_depth
            return max(self._calculate_json_depth(item, current_depth + 1) for item in data)
        else:
            return current_depth
    
    def _read_pdf_file(self, file_path: Path) -> str:
        """Read PDF files (requires PyPDF2 or similar)"""
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except ImportError:
            return "PDF reading requires PyPDF2 library. Please install it to read PDF files."
        except Exception as e:
            return f"Error reading PDF: {str(e)}"
    
    def _read_document_file(self, file_path: Path) -> str:
        """Read document files (requires python-docx)"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            return "Document reading requires python-docx library. Please install it to read Word documents."
        except Exception as e:
            return f"Error reading document: {str(e)}"
    
    def _read_spreadsheet_file(self, file_path: Path) -> str:
        """Read spreadsheet files"""
        try:
            import pandas as pd
            
            # Read all sheets
            excel_data = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            
            text = f"Spreadsheet: {file_path.name}\n\n"
            
            for sheet_name, df in excel_data.items():
                text += f"=== Sheet: {sheet_name} ===\n"
                text += df.to_string(index=False) + "\n\n"
            
            return text
        except ImportError:
            return "Spreadsheet reading requires pandas and openpyxl libraries."
        except Exception as e:
            return f"Error reading spreadsheet: {str(e)}"
    
    def _read_presentation_file(self, file_path: Path) -> str:
        """Read presentation files (PowerPoint)"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text = f"Presentation: {file_path.name}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"=== Slide {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
            
            return text
        except ImportError:
            return "Presentation reading requires python-pptx library. Please install it to read PowerPoint files."
        except Exception as e:
            return f"Error reading presentation: {str(e)}"
    
    def _read_image_file(self, file_path: Path) -> str:
        """Read image files (basic info + OCR if available)"""
        try:
            from PIL import Image
            
            with Image.open(file_path) as img:
                info_text = f"Image: {file_path.name}\n"
                info_text += f"Format: {img.format}\n"
                info_text += f"Size: {img.size}\n"
                info_text += f"Mode: {img.mode}\n\n"
                
                # Try OCR if available
                try:
                    import pytesseract
                    extracted_text = pytesseract.image_to_string(img)
                    if extracted_text.strip():
                        info_text += "Extracted Text:\n" + extracted_text
                    else:
                        info_text += "No text found in image."
                except ImportError:
                    info_text += "OCR not available. Install pytesseract for text extraction."
                except Exception as ocr_e:
                    info_text += f"OCR failed: {str(ocr_e)}"
                
                return info_text
                
        except ImportError:
            return "Image reading requires PIL/Pillow library."
        except Exception as e:
            return f"Error reading image: {str(e)}"
    
    def _read_archive_file(self, file_path: Path) -> str:
        """Read archive files (list contents)"""
        try:
            if file_path.suffix.lower() == '.zip':
                import zipfile
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    file_list = zip_file.namelist()
                    
                    text = f"ZIP Archive: {file_path.name}\n"
                    text += f"Files: {len(file_list)}\n\n"
                    text += "Contents:\n"
                    
                    for file_name in file_list[:50]:  # Limit to first 50 files
                        text += f"  {file_name}\n"
                    
                    if len(file_list) > 50:
                        text += f"  ... and {len(file_list) - 50} more files\n"
                    
                    return text
            else:
                return f"Archive type {file_path.suffix} not supported for content reading."
                
        except Exception as e:
            return f"Error reading archive: {str(e)}"
    
    def get_enhanced_analysis(self, file_id: str, user_id: str, enhanced_processor=None) -> Dict[str, Any]:
        """
        Get enhanced analysis of a file using the enhanced processor
        
        Args:
            file_id: File ID
            user_id: User ID
            enhanced_processor: Enhanced file processor instance
            
        Returns:
            Enhanced analysis results
        """
        try:
            if not enhanced_processor:
                return {"success": False, "error": "Enhanced processor not available"}
            
            # Get file info
            file_info = self.get_file_info(file_id, user_id)
            if not file_info:
                return {"success": False, "error": "File not found"}
            
            file_path = file_info["file_path"]
            file_type = file_info["file_type"]
            
            # Perform enhanced analysis
            return enhanced_processor.analyze_file(file_path, file_type, user_id)
            
        except Exception as e:
            self.logger.error(f"Enhanced analysis error: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def get_file_info(self, file_id: str, user_id: str) -> Optional[Dict[str, Any]]:
        """Get file metadata"""
        try:
            metadata_file = self.upload_dir / "file_metadata.json"
            if not metadata_file.exists():
                return None
            
            with open(metadata_file, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
            
            file_info = metadata.get(file_id)
            if file_info and file_info.get("user_id") == user_id:
                return file_info
            
            return None
            
        except Exception as e:
            self.logger.error(f"Error getting file info {file_id}: {str(e)}")
            return None
    
    def list_user_files(self, user_id: str) -> List[Dict[str, Any]]:
        """List files for a specific user"""
        try:
            metadata_file = self.upload_dir / "file_metadata.json"
            if not metadata_file.exists():
                return []
            
            with open(metadata_file, 'r', encoding='utf-8') as f:
                metadata = json.load(f)
            
            user_files = []
            for file_info in metadata.values():
                if file_info.get("user_id") == user_id:
                    # Don't include full file path in response
                    safe_info = {
                        "file_id": file_info["file_id"],
                        "original_name": file_info["original_name"],
                        "file_size": file_info["file_size"],
                        "file_type": file_info["file_type"],
                        "category": file_info["category"],
                        "upload_time": file_info["upload_time"]
                    }
                    user_files.append(safe_info)
            
            # Sort by upload time (newest first)
            user_files.sort(key=lambda x: x["upload_time"], reverse=True)
            return user_files
            
        except Exception as e:
            self.logger.error(f"Error listing files for user {user_id}: {str(e)}")
            return []
    
    def delete_file(self, file_id: str, user_id: str) -> bool:
        """Delete a file"""
        try:
            file_info = self.get_file_info(file_id, user_id)
            if not file_info:
                return False
            
            # Delete physical file
            file_path = Path(file_info["file_path"])
            if file_path.exists():
                file_path.unlink()
            
            # Remove from metadata
            metadata_file = self.upload_dir / "file_metadata.json"
            if metadata_file.exists():
                with open(metadata_file, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)
                
                if file_id in metadata:
                    del metadata[file_id]
                    
                    with open(metadata_file, 'w', encoding='utf-8') as f:
                        json.dump(metadata, f, indent=2, ensure_ascii=False)
            
            self.logger.info(f"File deleted: {file_id} for user {user_id}")
            return True
            
        except Exception as e:
            self.logger.error(f"Error deleting file {file_id}: {str(e)}")
            return False

def main():
    """Main function for testing file handler"""
    logging.basicConfig(level=logging.INFO)
    
    handler = FileHandler()
    print(f"File handler initialized with upload directory: {handler.upload_dir}")
    print(f"Supported file types: {list(handler.supported_types.keys())}")

if __name__ == "__main__":
    main()