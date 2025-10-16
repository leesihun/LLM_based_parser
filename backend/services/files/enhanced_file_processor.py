#!/usr/bin/env python3
"""
Enhanced File Processing Module
Advanced file analysis including PDF intelligence, code analysis, and multimodal processing
"""

import os
import json
import logging
import hashlib
import mimetypes
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import re
import ast
import zipfile
import tempfile

# Import libraries with fallbacks
try:
    import PyPDF2
    from pdfplumber import PDF as PDFPlumber
    HAS_ADVANCED_PDF = True
except ImportError:
    HAS_ADVANCED_PDF = False

try:
    from docx import Document
    from pptx import Presentation
    HAS_OFFICE_SUPPORT = True
except ImportError:
    HAS_OFFICE_SUPPORT = False

try:
    from PIL import Image, ImageEnhance
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    import pandas as pd
    import numpy as np
    HAS_DATA_ANALYSIS = True
except ImportError:
    HAS_DATA_ANALYSIS = False


class EnhancedFileProcessor:
    """Advanced file processing with intelligence features"""
    
    def __init__(self, config: Optional[Dict] = None):
        """
        Initialize enhanced file processor
        
        Args:
            config: Configuration dictionary
        """
        self.config = config or {}
        self.logger = logging.getLogger(__name__)
        
        # Enhanced file type support
        self.supported_types = {
            # Text files
            '.txt': 'text',
            '.md': 'text',
            '.rst': 'text',
            '.log': 'text',
            
            # Code files
            '.py': 'code',
            '.js': 'code',
            '.ts': 'code',
            '.jsx': 'code',
            '.tsx': 'code',
            '.java': 'code',
            '.cpp': 'code',
            '.c': 'code',
            '.h': 'code',
            '.cs': 'code',
            '.php': 'code',
            '.rb': 'code',
            '.go': 'code',
            '.rs': 'code',
            '.swift': 'code',
            '.kt': 'code',
            '.scala': 'code',
            '.sql': 'code',
            '.html': 'code',
            '.css': 'code',
            '.json': 'code',
            '.xml': 'code',
            '.yaml': 'code',
            '.yml': 'code',
            '.toml': 'code',
            '.ini': 'code',
            '.conf': 'code',
            
            # Documents
            '.pdf': 'document',
            '.docx': 'document',
            '.doc': 'document',
            '.pptx': 'presentation',
            '.ppt': 'presentation',
            '.xlsx': 'spreadsheet',
            '.xls': 'spreadsheet',
            '.csv': 'spreadsheet',
            
            # Images
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.bmp': 'image',
            '.tiff': 'image',
            '.tif': 'image',
            '.webp': 'image',
            
            # Archives
            '.zip': 'archive',
            '.tar': 'archive',
            '.gz': 'archive',
            '.rar': 'archive',
            '.7z': 'archive'
        }
        
        # Programming language keywords for analysis
        self.language_keywords = {
            '.py': ['def', 'class', 'import', 'from', 'if', 'for', 'while', 'try', 'except'],
            '.js': ['function', 'class', 'const', 'let', 'var', 'if', 'for', 'while', 'try', 'catch'],
            '.java': ['public', 'private', 'class', 'interface', 'extends', 'implements', 'if', 'for', 'while'],
            '.cpp': ['class', 'struct', 'namespace', 'template', 'if', 'for', 'while', 'try', 'catch'],
            '.cs': ['public', 'private', 'class', 'interface', 'namespace', 'if', 'for', 'while', 'try', 'catch']
        }
    
    def analyze_file(self, file_path: str, file_type: str, user_id: str) -> Dict[str, Any]:
        """
        Comprehensive file analysis
        
        Args:
            file_path: Path to file
            file_type: File extension
            user_id: User ID for context
            
        Returns:
            Analysis results dictionary
        """
        try:
            file_path_obj = Path(file_path)
            category = self.supported_types.get(file_type, 'unknown')
            
            analysis = {
                'success': True,
                'file_path': file_path,
                'file_type': file_type,
                'category': category,
                'size': file_path_obj.stat().st_size,
                'timestamp': datetime.now().isoformat(),
                'analysis': {}
            }
            
            # Route to appropriate analyzer
            if category == 'text':
                analysis['analysis'] = self._analyze_text_file(file_path_obj)
            elif category == 'code':
                analysis['analysis'] = self._analyze_code_file(file_path_obj, file_type)
            elif category == 'document':
                analysis['analysis'] = self._analyze_document_file(file_path_obj, file_type)
            elif category == 'presentation':
                analysis['analysis'] = self._analyze_presentation_file(file_path_obj)
            elif category == 'spreadsheet':
                analysis['analysis'] = self._analyze_spreadsheet_file(file_path_obj)
            elif category == 'image':
                analysis['analysis'] = self._analyze_image_file(file_path_obj)
            elif category == 'archive':
                analysis['analysis'] = self._analyze_archive_file(file_path_obj)
            else:
                analysis['analysis'] = self._analyze_generic_file(file_path_obj)
            
            return analysis
            
        except Exception as e:
            self.logger.error(f"Error analyzing file {file_path}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'file_path': file_path,
                'file_type': file_type
            }
    
    def _analyze_text_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze text files"""
        try:
            content = self._read_text_with_encoding(file_path)
            
            # Basic text statistics
            lines = content.split('\n')
            words = content.split()
            characters = len(content)
            
            # Language detection (simple heuristic)
            language = self._detect_text_language(content)
            
            # Extract key information
            headings = self._extract_headings(content)
            urls = self._extract_urls(content)
            emails = self._extract_emails(content)
            
            return {
                'content': content,
                'statistics': {
                    'lines': len(lines),
                    'words': len(words),
                    'characters': characters,
                    'non_empty_lines': len([line for line in lines if line.strip()])
                },
                'language': language,
                'extracted_info': {
                    'headings': headings,
                    'urls': urls,
                    'emails': emails
                },
                'type': 'text_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'text_analysis'}
    
    def _analyze_code_file(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Analyze code files with advanced metrics"""
        try:
            content = self._read_text_with_encoding(file_path)
            
            # Basic code statistics
            lines = content.split('\n')
            code_lines = [line for line in lines if line.strip() and not line.strip().startswith('#')]
            comment_lines = [line for line in lines if line.strip().startswith('#') or line.strip().startswith('//')]
            
            # Extract functions, classes, imports
            functions = []
            classes = []
            imports = []
            
            if file_type == '.py':
                analysis = self._analyze_python_code(content)
                functions = analysis.get('functions', [])
                classes = analysis.get('classes', [])
                imports = analysis.get('imports', [])
            elif file_type in ['.js', '.ts', '.jsx', '.tsx']:
                analysis = self._analyze_javascript_code(content)
                functions = analysis.get('functions', [])
                classes = analysis.get('classes', [])
                imports = analysis.get('imports', [])
            
            # Calculate complexity metrics
            complexity = self._calculate_code_complexity(content, file_type)
            
            # Security analysis (basic)
            security_issues = self._detect_security_issues(content, file_type)
            
            return {
                'content': content,
                'statistics': {
                    'total_lines': len(lines),
                    'code_lines': len(code_lines),
                    'comment_lines': len(comment_lines),
                    'blank_lines': len(lines) - len(code_lines) - len(comment_lines),
                    'functions_count': len(functions),
                    'classes_count': len(classes),
                    'imports_count': len(imports)
                },
                'extracted_elements': {
                    'functions': functions,
                    'classes': classes,
                    'imports': imports
                },
                'complexity': complexity,
                'security_analysis': security_issues,
                'type': 'code_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'code_analysis'}
    
    def _analyze_document_file(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Analyze document files (PDF, Word)"""
        try:
            if file_type == '.pdf':
                return self._analyze_pdf_advanced(file_path)
            elif file_type in ['.docx', '.doc']:
                return self._analyze_word_document(file_path)
            else:
                return {'error': 'Unsupported document type', 'type': 'document_analysis'}
                
        except Exception as e:
            return {'error': str(e), 'type': 'document_analysis'}
    
    def _analyze_pdf_advanced(self, file_path: Path) -> Dict[str, Any]:
        """Advanced PDF analysis with table and image extraction"""
        try:
            analysis = {
                'text_content': '',
                'pages': 0,
                'tables': [],
                'images': [],
                'metadata': {},
                'structure': [],
                'type': 'pdf_analysis'
            }
            
            # Basic text extraction
            if HAS_ADVANCED_PDF:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    analysis['pages'] = len(pdf.pages)
                    analysis['metadata'] = pdf.metadata or {}
                    
                    full_text = []
                    for i, page in enumerate(pdf.pages):
                        # Extract text
                        page_text = page.extract_text()
                        if page_text:
                            full_text.append(f"=== Page {i+1} ===\n{page_text}")
                        
                        # Extract tables
                        tables = page.extract_tables()
                        for j, table in enumerate(tables):
                            if table:
                                analysis['tables'].append({
                                    'page': i+1,
                                    'table_id': j+1,
                                    'data': table,
                                    'rows': len(table),
                                    'columns': len(table[0]) if table else 0
                                })
                    
                    analysis['text_content'] = '\n\n'.join(full_text)
            else:
                # Fallback to basic PyPDF2
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    analysis['pages'] = len(reader.pages)
                    analysis['metadata'] = reader.metadata or {}
                    
                    full_text = []
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            full_text.append(f"=== Page {i+1} ===\n{page_text}")
                    
                    analysis['text_content'] = '\n\n'.join(full_text)
            
            # Extract document structure
            analysis['structure'] = self._extract_document_structure(analysis['text_content'])
            
            # Document statistics
            analysis['statistics'] = {
                'pages': analysis['pages'],
                'characters': len(analysis['text_content']),
                'words': len(analysis['text_content'].split()),
                'tables_count': len(analysis['tables']),
                'images_count': len(analysis['images'])
            }
            
            return analysis
            
        except Exception as e:
            return {'error': str(e), 'type': 'pdf_analysis'}
    
    def _analyze_word_document(self, file_path: Path) -> Dict[str, Any]:
        """Analyze Word documents"""
        try:
            if not HAS_OFFICE_SUPPORT:
                return {'error': 'Office document support not available', 'type': 'document_analysis'}
            
            doc = Document(file_path)
            
            # Extract text content
            paragraphs = []
            tables = []
            
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    para = None
                    for p in doc.paragraphs:
                        if p._element == element:
                            para = p
                            break
                    if para and para.text.strip():
                        paragraphs.append({
                            'text': para.text,
                            'style': para.style.name if para.style else 'Normal'
                        })
                elif element.tag.endswith('tbl'):  # Table
                    table = None
                    for t in doc.tables:
                        if t._element == element:
                            table = t
                            break
                    if table:
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        tables.append({
                            'data': table_data,
                            'rows': len(table_data),
                            'columns': len(table_data[0]) if table_data else 0
                        })
            
            full_text = '\n'.join([p['text'] for p in paragraphs])
            
            return {
                'content': full_text,
                'paragraphs': paragraphs,
                'tables': tables,
                'statistics': {
                    'paragraphs': len(paragraphs),
                    'tables': len(tables),
                    'characters': len(full_text),
                    'words': len(full_text.split())
                },
                'structure': self._extract_document_structure(full_text),
                'type': 'word_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'document_analysis'}
    
    def _analyze_presentation_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze PowerPoint presentations"""
        try:
            if not HAS_OFFICE_SUPPORT:
                return {'error': 'Office document support not available', 'type': 'presentation_analysis'}
            
            prs = Presentation(file_path)
            
            slides_content = []
            total_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                shapes_info = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)
                        total_text.append(shape.text)
                    
                    shapes_info.append({
                        'type': str(type(shape).__name__),
                        'has_text': hasattr(shape, 'text'),
                        'text_length': len(shape.text) if hasattr(shape, 'text') else 0
                    })
                
                slides_content.append({
                    'slide_number': i + 1,
                    'text_content': '\n'.join(slide_text),
                    'shapes_count': len(slide.shapes),
                    'shapes_info': shapes_info
                })
            
            full_text = '\n\n'.join(total_text)
            
            return {
                'content': full_text,
                'slides': slides_content,
                'statistics': {
                    'slides_count': len(prs.slides),
                    'total_characters': len(full_text),
                    'total_words': len(full_text.split()),
                    'slides_with_text': len([s for s in slides_content if s['text_content'].strip()])
                },
                'type': 'presentation_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'presentation_analysis'}
    
    def _analyze_spreadsheet_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze spreadsheet files"""
        try:
            if not HAS_DATA_ANALYSIS:
                return {'error': 'Data analysis libraries not available', 'type': 'spreadsheet_analysis'}
            
            # Read all sheets
            excel_data = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            
            sheets_analysis = []
            summary_text = []
            
            for sheet_name, df in excel_data.items():
                # Basic statistics
                sheet_info = {
                    'name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns),
                    'columns_names': df.columns.tolist(),
                    'data_types': df.dtypes.to_dict(),
                    'null_counts': df.isnull().sum().to_dict(),
                    'summary': df.describe(include='all').to_dict() if not df.empty else {}
                }
                
                sheets_analysis.append(sheet_info)
                summary_text.append(f"=== Sheet: {sheet_name} ===\n{df.to_string(index=False)}")
            
            return {
                'content': '\n\n'.join(summary_text),
                'sheets': sheets_analysis,
                'statistics': {
                    'sheets_count': len(excel_data),
                    'total_rows': sum(len(df) for df in excel_data.values()),
                    'total_columns': sum(len(df.columns) for df in excel_data.values())
                },
                'type': 'spreadsheet_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'spreadsheet_analysis'}
    
    def _analyze_image_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze image files with OCR"""
        try:
            if not HAS_OCR:
                return {'error': 'OCR capabilities not available', 'type': 'image_analysis'}
            
            # Open and analyze image
            with Image.open(file_path) as img:
                # Basic image info
                analysis = {
                    'format': img.format,
                    'mode': img.mode,
                    'size': img.size,
                    'width': img.width,
                    'height': img.height,
                    'type': 'image_analysis'
                }
                
                # OCR text extraction
                try:
                    # Enhance image for better OCR
                    enhanced_img = ImageEnhance.Contrast(img).enhance(2)
                    enhanced_img = ImageEnhance.Sharpness(enhanced_img).enhance(2)
                    
                    # Extract text
                    extracted_text = pytesseract.image_to_string(enhanced_img)
                    analysis['extracted_text'] = extracted_text
                    analysis['text_confidence'] = self._calculate_ocr_confidence(extracted_text)
                    
                    # Extract data if it looks like a table/chart
                    if self._appears_to_be_table_or_chart(extracted_text):
                        analysis['data_extraction'] = self._extract_data_from_image_text(extracted_text)
                    
                except Exception as ocr_error:
                    analysis['ocr_error'] = str(ocr_error)
                    analysis['extracted_text'] = ''
                
                return analysis
                
        except Exception as e:
            return {'error': str(e), 'type': 'image_analysis'}
    
    def _analyze_archive_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze archive files"""
        try:
            if file_path.suffix.lower() == '.zip':
                return self._analyze_zip_file(file_path)
            else:
                return {'error': 'Archive type not supported yet', 'type': 'archive_analysis'}
                
        except Exception as e:
            return {'error': str(e), 'type': 'archive_analysis'}
    
    def _analyze_zip_file(self, file_path: Path) -> Dict[str, Any]:
        """Analyze ZIP archives"""
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                file_list = zip_file.filelist
                
                files_info = []
                total_size = 0
                file_types = {}
                
                for file_info in file_list:
                    if not file_info.is_dir():
                        file_ext = Path(file_info.filename).suffix.lower()
                        file_types[file_ext] = file_types.get(file_ext, 0) + 1
                        total_size += file_info.file_size
                        
                        files_info.append({
                            'filename': file_info.filename,
                            'size': file_info.file_size,
                            'compressed_size': file_info.compress_size,
                            'modified_date': str(file_info.date_time),
                            'file_type': file_ext
                        })
                
                # Check if this looks like a code project
                is_code_project = self._detect_code_project(files_info)
                
                analysis = {
                    'files_count': len([f for f in file_list if not f.is_dir()]),
                    'directories_count': len([f for f in file_list if f.is_dir()]),
                    'total_uncompressed_size': total_size,
                    'file_types': file_types,
                    'files': files_info[:50],  # Limit to first 50 files
                    'is_code_project': is_code_project,
                    'project_analysis': {},
                    'type': 'zip_analysis'
                }
                
                # If it's a code project, do deeper analysis
                if is_code_project:
                    analysis['project_analysis'] = self._analyze_code_project_structure(files_info)
                
                return analysis
                
        except Exception as e:
            return {'error': str(e), 'type': 'archive_analysis'}
    
    def _analyze_generic_file(self, file_path: Path) -> Dict[str, Any]:
        """Generic file analysis"""
        try:
            file_stat = file_path.stat()
            
            return {
                'size': file_stat.st_size,
                'modified': datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
                'created': datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
                'mime_type': mimetypes.guess_type(str(file_path))[0],
                'type': 'generic_analysis'
            }
            
        except Exception as e:
            return {'error': str(e), 'type': 'generic_analysis'}
    
    # Helper methods
    
    def _read_text_with_encoding(self, file_path: Path) -> str:
        """Read text file with proper encoding detection"""
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        # Last resort: read as binary and decode with error handling
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    
    def _detect_text_language(self, content: str) -> str:
        """Simple language detection"""
        # This is a very basic implementation
        # For production, consider using proper language detection libraries
        
        english_words = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by']
        korean_chars = re.findall(r'[가-힣]', content)
        chinese_chars = re.findall(r'[一-龯]', content)
        japanese_chars = re.findall(r'[ひらがな-ヿ]|[カタカナ-ヿ]', content)
        
        if len(korean_chars) > 10:
            return 'korean'
        elif len(chinese_chars) > 10:
            return 'chinese'
        elif len(japanese_chars) > 10:
            return 'japanese'
        elif any(word in content.lower() for word in english_words):
            return 'english'
        else:
            return 'unknown'
    
    def _extract_headings(self, content: str) -> List[str]:
        """Extract headings from text"""
        headings = []
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            # Markdown headings
            if line.startswith('#'):
                headings.append(line)
            # Lines that might be headings (short, followed by empty line)
            elif len(line) < 100 and line.isupper() and len(line) > 5:
                headings.append(line)
        
        return headings[:20]  # Limit to first 20
    
    def _extract_urls(self, content: str) -> List[str]:
        """Extract URLs from text"""
        url_pattern = r'https?://[^\s<>"{\}|\\^`\[\]]+'
        urls = re.findall(url_pattern, content)
        return list(set(urls))[:20]  # Unique URLs, limit to 20
    
    def _extract_emails(self, content: str) -> List[str]:
        """Extract email addresses from text"""
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        emails = re.findall(email_pattern, content)
        return list(set(emails))[:20]  # Unique emails, limit to 20
    
    def _analyze_python_code(self, content: str) -> Dict[str, Any]:
        """Analyze Python code using AST"""
        try:
            tree = ast.parse(content)
            
            functions = []
            classes = []
            imports = []
            
            for node in ast.walk(tree):
                if isinstance(node, ast.FunctionDef):
                    functions.append({
                        'name': node.name,
                        'line': node.lineno,
                        'args': [arg.arg for arg in node.args.args],
                        'decorators': [ast.unparse(d) for d in node.decorator_list] if hasattr(ast, 'unparse') else []
                    })
                elif isinstance(node, ast.ClassDef):
                    classes.append({
                        'name': node.name,
                        'line': node.lineno,
                        'bases': [ast.unparse(base) for base in node.bases] if hasattr(ast, 'unparse') else [],
                        'methods': [n.name for n in node.body if isinstance(n, ast.FunctionDef)]
                    })
                elif isinstance(node, ast.Import):
                    for alias in node.names:
                        imports.append({
                            'type': 'import',
                            'module': alias.name,
                            'alias': alias.asname,
                            'line': node.lineno
                        })
                elif isinstance(node, ast.ImportFrom):
                    for alias in node.names:
                        imports.append({
                            'type': 'from',
                            'module': node.module,
                            'name': alias.name,
                            'alias': alias.asname,
                            'line': node.lineno
                        })
            
            return {
                'functions': functions,
                'classes': classes,
                'imports': imports
            }
            
        except SyntaxError:
            return {'functions': [], 'classes': [], 'imports': [], 'syntax_error': True}
        except Exception as e:
            return {'functions': [], 'classes': [], 'imports': [], 'error': str(e)}
    
    def _analyze_javascript_code(self, content: str) -> Dict[str, Any]:
        """Basic JavaScript code analysis using regex"""
        functions = []
        classes = []
        imports = []
        
        # Extract function declarations
        func_patterns = [
            r'function\s+(\w+)\s*\([^)]*\)',
            r'const\s+(\w+)\s*=\s*\([^)]*\)\s*=>',
            r'let\s+(\w+)\s*=\s*\([^)]*\)\s*=>',
            r'var\s+(\w+)\s*=\s*\([^)]*\)\s*=>',
            r'(\w+):\s*function\s*\([^)]*\)'
        ]
        
        for pattern in func_patterns:
            matches = re.finditer(pattern, content, re.MULTILINE)
            for match in matches:
                functions.append({
                    'name': match.group(1),
                    'line': content[:match.start()].count('\n') + 1
                })
        
        # Extract class declarations
        class_matches = re.finditer(r'class\s+(\w+)', content, re.MULTILINE)
        for match in class_matches:
            classes.append({
                'name': match.group(1),
                'line': content[:match.start()].count('\n') + 1
            })
        
        # Extract imports
        import_patterns = [
            r'import\s+{([^}]+)}\s+from\s+[\'"]([^\'"]+)[\'"]',
            r'import\s+(\w+)\s+from\s+[\'"]([^\'"]+)[\'"]',
            r'const\s+(\w+)\s*=\s*require\([\'"]([^\'"]+)[\'"]\)'
        ]
        
        for pattern in import_patterns:
            matches = re.finditer(pattern, content, re.MULTILINE)
            for match in matches:
                imports.append({
                    'type': 'import',
                    'module': match.group(2) if len(match.groups()) > 1 else match.group(1),
                    'line': content[:match.start()].count('\n') + 1
                })
        
        return {
            'functions': functions,
            'classes': classes,
            'imports': imports
        }
    
    def _calculate_code_complexity(self, content: str, file_type: str) -> Dict[str, Any]:
        """Calculate basic code complexity metrics"""
        lines = content.split('\n')
        
        # Basic metrics
        total_lines = len(lines)
        code_lines = len([line for line in lines if line.strip() and not line.strip().startswith('#')])
        
        # Cyclomatic complexity (very basic approximation)
        complexity_keywords = ['if', 'for', 'while', 'case', 'catch', 'else', '&&', '||']
        complexity_count = 0
        
        for line in lines:
            for keyword in complexity_keywords:
                complexity_count += line.lower().count(keyword)
        
        # Nesting level (approximate)
        max_nesting = 0
        current_nesting = 0
        
        for line in lines:
            stripped = line.lstrip()
            if stripped:
                # Count indentation
                indent = len(line) - len(stripped)
                if file_type == '.py':
                    current_nesting = indent // 4
                else:
                    current_nesting = line.count('{') - line.count('}')
                max_nesting = max(max_nesting, current_nesting)
        
        return {
            'total_lines': total_lines,
            'code_lines': code_lines,
            'complexity_score': complexity_count,
            'max_nesting_level': max_nesting,
            'maintainability_index': min(100, max(0, 171 - 5.2 * np.log(max(1, code_lines)) - 0.23 * complexity_count - 16.2 * np.log(max(1, max_nesting)))) if HAS_DATA_ANALYSIS else 'N/A'
        }
    
    def _detect_security_issues(self, content: str, file_type: str) -> List[Dict[str, Any]]:
        """Detect potential security issues in code"""
        issues = []
        lines = content.split('\n')
        
        # Common security anti-patterns
        security_patterns = {
            'sql_injection': [r'execute\s*\(\s*[\'"].*%.*[\'"]', r'query\s*\(\s*[\'"].*\+.*[\'"]'],
            'hardcoded_secrets': [r'password\s*=\s*[\'"][^\'"]+[\'"]', r'api_key\s*=\s*[\'"][^\'"]+[\'"]', r'secret\s*=\s*[\'"][^\'"]+[\'"]'],
            'unsafe_eval': [r'eval\s*\(', r'exec\s*\('],
            'dangerous_imports': [r'import\s+os', r'import\s+subprocess', r'from\s+os\s+import']
        }
        
        for i, line in enumerate(lines, 1):
            for issue_type, patterns in security_patterns.items():
                for pattern in patterns:
                    if re.search(pattern, line.lower()):
                        issues.append({
                            'type': issue_type,
                            'line': i,
                            'code': line.strip(),
                            'severity': 'medium'
                        })
        
        return issues[:10]  # Limit to first 10 issues
    
    def _extract_document_structure(self, content: str) -> List[Dict[str, Any]]:
        """Extract document structure (headings, sections)"""
        structure = []
        lines = content.split('\n')
        
        for i, line in enumerate(lines):
            line = line.strip()
            if line:
                # Markdown headings
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    structure.append({
                        'type': 'heading',
                        'level': level,
                        'text': line.lstrip('#').strip(),
                        'line': i + 1
                    })
                # Potential section headers (all caps, short)
                elif line.isupper() and len(line) < 50 and len(line) > 3:
                    structure.append({
                        'type': 'section',
                        'text': line,
                        'line': i + 1
                    })
        
        return structure[:20]  # Limit to first 20 elements
    
    def _calculate_ocr_confidence(self, text: str) -> float:
        """Calculate OCR confidence based on text characteristics"""
        if not text.strip():
            return 0.0
        
        # Basic heuristics for OCR quality
        score = 0.0
        
        # Check for common OCR errors
        weird_chars = len(re.findall(r'[^\w\s.,!?;:\'"()-]', text))
        total_chars = len(text)
        
        if total_chars > 0:
            score = max(0.0, 1.0 - (weird_chars / total_chars))
        
        # Boost score for recognizable words
        words = text.split()
        if words:
            common_words = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by']
            common_word_count = sum(1 for word in words if word.lower() in common_words)
            score += (common_word_count / len(words)) * 0.2
        
        return min(1.0, score)
    
    def _appears_to_be_table_or_chart(self, text: str) -> bool:
        """Check if extracted text appears to be tabular data or chart"""
        lines = text.split('\n')
        
        # Look for table-like patterns
        numeric_lines = 0
        for line in lines:
            if re.search(r'\d+', line) and ('|' in line or '\t' in line or len(re.findall(r'\s{2,}', line)) > 2):
                numeric_lines += 1
        
        return numeric_lines > 3
    
    def _extract_data_from_image_text(self, text: str) -> Dict[str, Any]:
        """Extract structured data from OCR text"""
        # This is a basic implementation
        # For production, consider more sophisticated parsing
        
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        
        # Look for table-like structures
        rows = []
        for line in lines:
            # Split on multiple spaces or tabs
            cells = re.split(r'\s{2,}|\t', line)
            if len(cells) > 1:
                rows.append(cells)
        
        return {
            'rows': rows[:10],  # Limit to first 10 rows
            'extracted': True if rows else False,
            'row_count': len(rows)
        }
    
    def _detect_code_project(self, files_info: List[Dict[str, Any]]) -> bool:
        """Detect if ZIP contains a code project"""
        code_extensions = ['.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go']
        config_files = ['package.json', 'requirements.txt', 'pom.xml', 'Makefile', 'setup.py', '.gitignore']
        
        code_file_count = sum(1 for f in files_info if f['file_type'] in code_extensions)
        config_file_count = sum(1 for f in files_info if any(cf in f['filename'] for cf in config_files))
        
        return code_file_count > 3 or config_file_count > 0
    
    def _analyze_code_project_structure(self, files_info: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Analyze code project structure"""
        languages = {}
        directories = set()
        
        for file_info in files_info:
            # Count languages
            ext = file_info['file_type']
            if ext:
                languages[ext] = languages.get(ext, 0) + 1
            
            # Track directory structure
            path_parts = Path(file_info['filename']).parts
            if len(path_parts) > 1:
                directories.add(path_parts[0])
        
        # Detect project type
        project_type = 'unknown'
        if any('package.json' in f['filename'] for f in files_info):
            project_type = 'javascript/nodejs'
        elif any('requirements.txt' in f['filename'] or 'setup.py' in f['filename'] for f in files_info):
            project_type = 'python'
        elif any('pom.xml' in f['filename'] for f in files_info):
            project_type = 'java/maven'
        elif any('.csproj' in f['filename'] for f in files_info):
            project_type = 'csharp/.net'
        
        return {
            'project_type': project_type,
            'languages': languages,
            'main_directories': list(directories),
            'file_count': len(files_info),
            'estimated_complexity': 'high' if len(files_info) > 100 else 'medium' if len(files_info) > 20 else 'low'
        }


def main():
    """Main function for testing"""
    processor = EnhancedFileProcessor()
    print(f"Enhanced File Processor initialized")
    print(f"Supported file types: {list(processor.supported_types.keys())}")
    print(f"OCR available: {HAS_OCR}")
    print(f"Advanced PDF available: {HAS_ADVANCED_PDF}")
    print(f"Office support available: {HAS_OFFICE_SUPPORT}")


if __name__ == "__main__":
    main()